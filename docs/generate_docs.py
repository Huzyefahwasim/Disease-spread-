from docx import Document
from pptx import Presentation

def create_docx():
    doc = Document()
    doc.add_heading('IT3016 Simulation and Modelling - Project Report', 0)
    doc.add_heading('Disease Spread Simulator', 1)
    
    doc.add_heading('1. Introduction', level=2)
    doc.add_paragraph('This project simulates the spread of an infectious disease using an SIR (Susceptible-Infectious-Recovered) model. It aims to fulfill the requirements of CLO.3 and CLO.4 by providing a complete simulation environment coupled with Machine Learning predictions.')
    
    doc.add_heading('2. Dataset Analysis', level=2)
    doc.add_paragraph('The dataset was synthetically generated using an SIR model with added Gaussian noise to mimic real-world fluctuations in reporting. Exploratory Data Analysis (EDA) was performed to understand the relationships between stringency indices, testing rates, and new case counts.')
    
    doc.add_heading('3. Simulation Model', level=2)
    doc.add_paragraph('The core simulation uses the SIR compartmental model:')
    doc.add_paragraph('Susceptible (S): Individuals who can contract the disease.', style='List Bullet')
    doc.add_paragraph('Infectious (I): Individuals who have the disease and can transmit it.', style='List Bullet')
    doc.add_paragraph('Recovered (R): Individuals who have recovered and are immune.', style='List Bullet')
    doc.add_paragraph('Parameters such as transmission rate and recovery rate can be interactively modified in the Streamlit dashboard.')
    
    doc.add_heading('4. Machine Learning Models', level=2)
    doc.add_paragraph('Three models were trained to predict future infection rates based on a 7-day lag history and external factors (Stringency Index and Testing Rate):')
    doc.add_paragraph('Linear Regression: Serves as a baseline model.', style='List Number')
    doc.add_paragraph('Random Forest Regressor: Captures non-linear relationships.', style='List Number')
    doc.add_paragraph('Gradient Boosting Regressor: Provides robust predictions through an ensemble approach.', style='List Number')
    
    doc.add_heading('5. Conclusion', level=2)
    doc.add_paragraph('The Streamlit dashboard successfully integrates the simulation, dataset analysis, and ML predictions into a single interactive platform, demonstrating a comprehensive understanding of simulation and modelling techniques.')
    
    doc.save('docs/Project_Report.docx')
    print("Created docs/Project_Report.docx")

def create_pptx():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Disease Spread Simulator"
    subtitle.text = "IT3016 Simulation and Modelling\nSemester Project\nBy: [Your Name]"
    
    # Slides
    slides_content = [
        ("Introduction & Problem Statement", [
            "Understanding the rapid spread of infectious diseases.",
            "The need for computational models to predict and mitigate outbreaks.",
            "How simulation helps policymakers in decision making."
        ]),
        ("Project Objectives", [
            "Fulfill requirements for CLO.3 and CLO.4.",
            "Build a mathematical simulation of disease dynamics.",
            "Integrate Machine Learning to predict future trends based on historical data."
        ]),
        ("Theoretical Background: The SIR Model", [
            "Susceptible (S): Population at risk of contracting the disease.",
            "Infectious (I): Population currently infected and contagious.",
            "Recovered (R): Population that has recovered and gained immunity.",
            "Governed by transmission rate (beta) and recovery rate (gamma)."
        ]),
        ("Dataset Generation & EDA", [
            "Generating synthetic data using an SIR model with Gaussian noise.",
            "Features included: Stringency Index and Testing Rates.",
            "Exploratory Data Analysis: Correlation matrices and time-series plotting."
        ]),
        ("Machine Learning Architecture", [
            "Time-series forecasting using 7-day lag features.",
            "Model 1: Linear Regression (Baseline model).",
            "Model 2: Random Forest Regressor (Handles non-linear relationships).",
            "Model 3: Gradient Boosting Regressor (Ensemble learning for high accuracy)."
        ]),
        ("Model Evaluation & Comparison", [
            "Evaluation metrics: Mean Squared Error (MSE) and R-squared (R2).",
            "Random Forest generally performs best on non-linear epidemic curves.",
            "Discussion on overfitting vs. generalization."
        ]),
        ("Streamlit Dashboard Architecture", [
            "Component 1: EDA Viewer (Dataframes, Matplotlib plots).",
            "Component 2: Interactive Simulator (Adjustable transmission/recovery).",
            "Component 3: ML Predictor Interface for future cases."
        ]),
        ("Interactive Simulation & Results", [
            "Demonstration of the interactive simulation.",
            "Visualizing 'flattening the curve' by lowering transmission rates.",
            "Real-time updates dynamically driven by Streamlit."
        ]),
        ("Challenges & Limitations", [
            "Assumptions of standard SIR (constant population, homogenous mixing).",
            "Limitations of synthetic datasets compared to real-world noisy data.",
            "Time-series forecasting limits using simple lag features."
        ]),
        ("Conclusion & Future Enhancements", [
            "Successfully created a full-stack data science & simulation application.",
            "Future work: Expanding to SEIR models (adding 'Exposed' compartment).",
            "Adding deep learning (LSTMs) for more robust predictions."
        ]),
        ("Q&A", [
            "Thank you for your attention!",
            "Are there any questions?"
        ])
    ]
    
    bullet_slide_layout = prs.slide_layouts[1]
    for title_text, bullets in slides_content:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        tf = body_shape.text_frame
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
    prs.save('docs/Presentation_Slides.pptx')
    print("Created docs/Presentation_Slides.pptx")

if __name__ == '__main__':
    create_docx()
    create_pptx()
